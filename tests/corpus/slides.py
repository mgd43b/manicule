"""Generated `.pptx` fixtures: typical, structurally hard, degenerate and hostile.

Generated rather than committed (``docs/parsing.md`` §3.5), so each deck's structure is
reviewable as code and the hostile packages exist without being stored.

The structurally hard deck is built around the three ways shape geometry stops being a usable
box, because each one silently costs a rectangle and leaves a citation that still names the
right slide: **overlapping** shapes (which must keep their own boxes), two shapes at
**identical** coordinates (which cannot be told apart by a box and must go page-level), and a
shape whose ``<a:xfrm>`` has been **removed** (which reports no position at all). The last one
is produced by editing the saved package: python-pptx's default layouts all carry geometry, so
there is no way to ask it for a shape without any.
"""

from __future__ import annotations

import base64
import io
import zipfile
from collections.abc import Callable
from pathlib import Path

import pptx
from pptx.shapes.autoshape import Shape
from pptx.shapes.graphfrm import GraphicFrame
from pptx.slide import Slide
from pptx.util import Inches, Length, Pt

__all__ = ["build"]

_TITLE_AND_CONTENT = 1
_BLANK = 6

_PNG_1X1 = base64.b64decode(
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAYAAAAfFcSJAAAADUlEQVR42mP8z8AAAAMBAQDJ/pLvAAAAAElFTkSuQmCC"
)
"""A one-pixel PNG. A slide holding only this has no text, and no text is the honest answer."""

_XFRM_OPEN = b"<a:xfrm>"
_XFRM_CLOSE = b"</a:xfrm>"


def build(dest: Path) -> None:
    """Write this format's fixtures into ``dest``."""
    dest.mkdir(parents=True, exist_ok=True)
    _typical(dest / "slides_typical.pptx")
    _structurally_hard(dest / "slides_structurally_hard.pptx")
    _image_only(dest / "slides_degenerate_image_only.pptx")
    _no_slides(dest / "slides_degenerate_no_slides.pptx")
    _blank_slide(dest / "slides_degenerate_blank_slide.pptx")
    (dest / "slides_degenerate_zero_bytes.pptx").write_bytes(b"")
    _astral(dest / "slides_hostile_astral.pptx")
    _large(dest / "slides-large.pptx")
    _truncated(dest / "slides_hostile_truncated.pptx")
    _plain_zip(dest / "slides_hostile_plain_zip.pptx")


def _typical(path: Path) -> None:
    presentation = pptx.Presentation()

    first = presentation.slides.add_slide(presentation.slide_layouts[_TITLE_AND_CONTENT])
    _set_title(first, "Quarterly platform review")
    _set_body(
        first,
        [
            ("Latency held inside the objective for eleven weeks", 0),
            ("Two regions carried the growth", 1),
        ],
    )
    _set_notes(
        first, "Mention that the eleven-week figure excludes the maintenance window in April."
    )

    second = presentation.slides.add_slide(presentation.slide_layouts[_TITLE_AND_CONTENT])
    _set_title(second, "Where the budget went")
    _set_body(second, [("Storage grew faster than compute, by roughly a third", 0)])

    third = presentation.slides.add_slide(presentation.slide_layouts[_TITLE_AND_CONTENT])
    _set_title(third, "Asks for next quarter")
    _set_body(third, [("One more reviewer on the release rota", 0)])

    presentation.save(str(path))


def _structurally_hard(path: Path) -> None:
    presentation = pptx.Presentation()

    overlapping = presentation.slides.add_slide(presentation.slide_layouts[_BLANK])
    _add_textbox(overlapping, "Backing panel text behind the callout", Inches(1), Inches(1))
    # Overlapping, not coincident: each shape keeps its own box, and the boxes intersect.
    _add_textbox(overlapping, "Callout text drawn over the panel", Inches(2), Inches(1.5))

    coincident = presentation.slides.add_slide(presentation.slide_layouts[_BLANK])
    _add_textbox(coincident, "First stacked shape, same coordinates", Inches(1), Inches(2))
    _add_textbox(coincident, "Second stacked shape, same coordinates", Inches(1), Inches(2))

    deep = presentation.slides.add_slide(presentation.slide_layouts[_TITLE_AND_CONTENT])
    _set_title(deep, "Nesting five levels down")
    _set_body(
        deep,
        [(f"Indent {level} of five, {'sub' * (level + 1)}ordinate", level) for level in range(5)],
    )

    tabular = presentation.slides.add_slide(presentation.slide_layouts[_BLANK])
    frame = tabular.shapes.add_table(3, 3, Inches(0.5), Inches(0.5), Inches(9), Inches(2))
    frame.table.first_row = True
    _fill_table(
        frame,
        [
            ["Region", "Requests", "Errors"],
            ["EMEA", "48 200", "12"],
            ["APAC", "31 750", "31"],
        ],
    )

    ungeometried = presentation.slides.add_slide(presentation.slide_layouts[_BLANK])
    _add_textbox(ungeometried, "This shape reports no position at all", Inches(1), Inches(1))

    buffer = io.BytesIO()
    presentation.save(buffer)
    # The last slide is rewritten with its geometry removed. Editing the package is the only
    # route: every layout python-pptx ships carries a position, so a placeholder inheriting
    # nothing cannot be constructed through the API.
    path.write_bytes(_strip_geometry(buffer.getvalue(), "ppt/slides/slide5.xml"))


def _image_only(path: Path) -> None:
    """A slide whose only shape is a picture. No text, so no blocks — not a failure."""
    presentation = pptx.Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[_BLANK])
    slide.shapes.add_picture(io.BytesIO(_PNG_1X1), Inches(2), Inches(2), Inches(1), Inches(1))
    presentation.save(str(path))


def _no_slides(path: Path) -> None:
    """A presentation with no slides at all."""
    pptx.Presentation().save(str(path))


def _blank_slide(path: Path) -> None:
    """One slide with no shapes on it."""
    presentation = pptx.Presentation()
    presentation.slides.add_slide(presentation.slide_layouts[_BLANK])
    presentation.save(str(path))


def _astral(path: Path) -> None:
    presentation = pptx.Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[_TITLE_AND_CONTENT])
    # Ambiguous-character warnings are suppressed on purpose: astral codepoints in a title are
    # the fixture, and a citation has to reproduce them exactly.
    _set_title(slide, "図表 𝔊raphs")  # noqa: RUF001
    _set_body(slide, [("Supplementary plane glyphs: 🜃 𠀋 and 𝕐", 0)])  # noqa: RUF001
    presentation.save(str(path))


def _large(path: Path) -> None:
    """A deck long enough to exercise the streaming path."""
    presentation = pptx.Presentation()
    for number in range(1, 26):
        slide = presentation.slides.add_slide(presentation.slide_layouts[_TITLE_AND_CONTENT])
        _set_title(slide, f"Week {number:02d} operations summary")
        _set_body(
            slide,
            [
                (
                    f"Week {number:02d}: {number * 137} deployments, {number} rollbacks, and a "
                    f"p99 of {200 + number} milliseconds across every region",
                    0,
                )
            ],
        )
    presentation.save(str(path))


def _set_notes(slide: Slide, text: str) -> None:
    """Set a slide's speaker notes, which are content on the slide they belong to."""
    frame = slide.notes_slide.notes_text_frame
    if frame is None:  # pragma: no cover - a notes slide always has a notes placeholder
        msg = "the notes slide has no notes placeholder"
        raise ValueError(msg)
    frame.text = text


def _set_title(slide: Slide, text: str) -> None:
    """Set a slide's title placeholder, which becomes the heading block and the path."""
    title = slide.shapes.title
    if title is None:  # pragma: no cover - the chosen layouts all have one
        msg = "the layout has no title placeholder"
        raise ValueError(msg)
    title.text_frame.text = text


def _set_body(slide: Slide, lines: list[tuple[str, int]]) -> None:
    """Fill the content placeholder with ``(text, indent level)`` paragraphs."""
    body = slide.placeholders[1]
    if not isinstance(body, Shape):  # pragma: no cover - the content placeholder is a Shape
        msg = f"the content placeholder holds no text frame: {type(body).__name__}"
        raise TypeError(msg)
    frame = body.text_frame
    first_text, first_level = lines[0]
    frame.text = first_text
    frame.paragraphs[0].level = first_level
    for text, level in lines[1:]:
        paragraph = frame.add_paragraph()
        paragraph.text = text
        paragraph.level = level


def _add_textbox(slide: Slide, text: str, left: Length, top: Length) -> None:
    """Add a positioned textbox of a fixed size."""
    box = slide.shapes.add_textbox(left, top, Inches(4), Inches(1.2))
    frame = box.text_frame
    frame.text = text
    frame.paragraphs[0].font.size = Pt(18)


def _fill_table(frame: GraphicFrame, rows: list[list[str]]) -> None:
    """Write a grid into a table graphic frame."""
    table = frame.table
    for row_index, row in enumerate(rows):
        for column_index, value in enumerate(row):
            table.cell(row_index, column_index).text = value


def _truncated(path: Path) -> None:
    """A package whose ``ppt/presentation.xml`` stops mid-tag. The parser must decline."""
    presentation = pptx.Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[_TITLE_AND_CONTENT])
    _set_title(slide, "Intact title")
    buffer = io.BytesIO()
    presentation.save(buffer)
    path.write_bytes(_truncate_member(buffer.getvalue(), "ppt/presentation.xml"))


def _plain_zip(path: Path) -> None:
    """A zip that is not an OOXML package at all, under a `.pptx` name."""
    buffer = io.BytesIO()
    with zipfile.ZipFile(buffer, "w", zipfile.ZIP_DEFLATED) as archive:
        archive.writestr("slides.txt", "not a presentation")
    path.write_bytes(buffer.getvalue())


def _rewrite(package: bytes, member: str, transform: Callable[[bytes], bytes]) -> bytes:
    """Rewrite one member of a zip, leaving the container and every other member valid."""
    source = zipfile.ZipFile(io.BytesIO(package))
    out = io.BytesIO()
    with zipfile.ZipFile(out, "w", zipfile.ZIP_DEFLATED) as archive:
        for info in source.infolist():
            blob = source.read(info.filename)
            archive.writestr(info.filename, transform(blob) if info.filename == member else blob)
    return out.getvalue()


def _truncate_member(package: bytes, member: str) -> bytes:
    return _rewrite(package, member, lambda blob: blob[: len(blob) // 2])


def _strip_geometry(package: bytes, member: str) -> bytes:
    """Remove every ``<a:xfrm>`` element from one slide part."""

    def transform(blob: bytes) -> bytes:
        while True:
            start = blob.find(_XFRM_OPEN)
            end = blob.find(_XFRM_CLOSE)
            if start < 0 or end < start:
                return blob
            blob = blob[:start] + blob[end + len(_XFRM_CLOSE) :]

    return _rewrite(package, member, transform)
