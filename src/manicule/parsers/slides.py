"""PPTX: a slide has a number and its shapes have boxes, so a citation has both.

``PageAnchor.page`` is the slide's **1-based position in presentation order**, because that
is what a viewer shows in the corner and what a person says out loud. It is deliberately not
a stable identifier: reordering a deck changes the document's ``version_token``, so the deck
is re-parsed and the anchors are rebuilt. The stable per-slide identifier is recorded in
``metadata.slide_id``, which is what lets a diff between two versions tell a moved slide from
a new one (``docs/parsing.md`` §2.4).

``rects`` comes from shape geometry, and three details in that conversion each produce a
wrong highlight rather than an error:

**Units and origin.** python-pptx reports EMU with a **top-left** origin, which is already
the convention :class:`~manicule.core.anchors.Rect` uses, so there is no y-flip here — unlike
the PDF path, where there is. The coordinates are divided by the slide's own width and height,
because ``Rect`` is normalised ``0.0``-``1.0``.

**Geometry that is not reported.** A placeholder inheriting its position from a layout that
declares none, or a shape whose ``<a:xfrm>`` was dropped by an editor, reports ``None`` for
``left``/``top``. That contributes ``rects=[]`` — an honest page-level anchor — rather than a
box at the origin, which would highlight the top-left corner of a slide the quotation is not
in. This is what the 0.20 page-level budget in ``docs/parsing.md`` §3.4 pays for.

**Geometry that is not unique.** Two shapes stacked at identical coordinates — a real habit
in decks where one shape is a backing panel for another — cannot be told apart by their box.
Such shapes go page-level too: a rectangle that covers two shapes' text claims text that was
not quoted, which is the merged-envelope mistake rule 3 of §2.1 forbids, arriving by a
different route.

Speaker notes are content and are emitted as ``prose`` on the slide they belong to. Their
anchor is page-level by construction: notes live on a notes *page* whose geometry is not the
slide's, so a box from it would point at a coordinate on a different sheet of paper.
"""

from __future__ import annotations

import io
import zipfile
from collections import Counter
from collections.abc import AsyncIterator, Iterable, Iterator, Sequence
from dataclasses import dataclass, field, replace

import pptx
from pptx.exc import PackageNotFoundError
from pptx.presentation import Presentation
from pptx.shapes.autoshape import Shape
from pptx.shapes.base import BaseShape
from pptx.shapes.graphfrm import GraphicFrame
from pptx.shapes.group import GroupShape
from pptx.slide import Slide
from pptx.table import Table
from pydantic import BaseModel, Field

from manicule.core.anchors import Anchor, PageAnchor, Rect
from manicule.core.content import BlockKind, Metadata, ParsedBlock, RawDocument
from manicule.core.errors import ParseError
from manicule.parsers.base import ParserProfile

__all__ = ["MEDIA_TYPE", "SlidesConfig", "SlidesParser"]

MEDIA_TYPE = "application/vnd.openxmlformats-officedocument.presentationml.presentation"

_CELL_SEPARATOR = "\t"
_ROW_SEPARATOR = "\n"

_COORDINATE_TOLERANCE = 1e-9
"""How close two normalised coordinates must be to name the same shape.

:meth:`SlidesParser.resolve` recomputes a shape's box with the same arithmetic
:meth:`SlidesParser.parse` used, so the values agree exactly; the tolerance is here so that a
last-bit difference in a stored-and-reloaded float cannot make an anchor stop resolving. It is
three orders of magnitude tighter than :data:`~manicule.core.anchors.EDGE_TOLERANCE`, which
absorbs transform noise — this one only absorbs round-tripping.
"""


class SlidesConfig(BaseModel):
    """Configuration for :class:`SlidesParser`."""

    include_speaker_notes: bool = Field(
        default=True,
        description="Index speaker notes as prose on their slide. On by default because the "
        "notes frequently carry the sentence the slide only gestures at; off for decks whose "
        "notes are a rehearsal script nobody should retrieve.",
    )


@dataclass(frozen=True, slots=True)
class _SlideItem:
    """One block-to-be from a slide, with the box it occupies if it reports one."""

    kind: BlockKind
    text: str
    rect: Rect | None
    metadata: Metadata = field(default_factory=Metadata)


def _emu(value: object) -> int | None:
    """A geometry value as a count of EMU, or ``None`` when the file states none.

    Typed as ``object`` on purpose. python-pptx annotates ``BaseShape.left`` as ``Length``
    while returning ``None`` for a shape with no ``<a:xfrm>``, and leaves
    ``Presentation.slide_width`` unannotated altogether — so the widening happens once, here,
    where "not a number of EMU" and "absent" collapse into the same honest answer.
    """
    return int(value) if isinstance(value, int) and not isinstance(value, bool) else None


def _slide_size(presentation: Presentation) -> tuple[int, int] | None:
    """The slide's EMU dimensions, or ``None`` when the file does not state them.

    Without them a fraction of the slide cannot be computed, so every anchor in the deck is
    page-level. Returning ``None`` rather than assuming 4:3 keeps that visible in the
    page-level ratio instead of producing boxes measured against a guess.
    """
    width = _emu(presentation.slide_width)  # pyright: ignore[reportUnknownMemberType,reportUnknownArgumentType] - python-pptx leaves slide_width unannotated
    height = _emu(presentation.slide_height)
    if width is None or height is None or width <= 0 or height <= 0:
        return None
    return width, height


def _shape_rect(shape: BaseShape, size: tuple[int, int] | None) -> Rect | None:
    """The shape's box as a fraction of the slide, or ``None`` if it has no usable one."""
    if size is None:
        return None
    left, top = _emu(shape.left), _emu(shape.top)
    width, height = _emu(shape.width), _emu(shape.height)
    if left is None or top is None or width is None or height is None:
        return None
    slide_width, slide_height = size
    # Clamped, not rejected: a shape hanging over the edge is visible in the part that is on
    # the slide, and that part is where the quotation is. A box outside 0.0-1.0 would be
    # refused by Rect, which would lose the anchor entirely.
    x0 = min(max(left / slide_width, 0.0), 1.0)
    y0 = min(max(top / slide_height, 0.0), 1.0)
    x1 = min(max((left + width) / slide_width, 0.0), 1.0)
    y1 = min(max((top + height) / slide_height, 0.0), 1.0)
    if x1 <= x0 or y1 <= y0:
        # Entirely off-slide, or zero-sized. There is no region a reader could be shown.
        return None
    return Rect(x0=x0, y0=y0, x1=x1, y1=y1)


def _render_table(table: Table) -> list[str]:
    """A slide table as one tab-separated line per row."""
    return [_CELL_SEPARATOR.join(cell.text.strip() for cell in row.cells) for row in table.rows]


def _table_metadata(table: Table, rows: Sequence[str]) -> Metadata:
    """What the chunker needs to split this table without losing its header.

    ``header_rows`` comes from the table's own first-row flag — the one PowerPoint records —
    and is 0 when the file does not set it. Nothing here reads formatting to decide.

    The row and column counts are taken from the rendered lines rather than from the table's
    own collections, so they describe exactly the text the chunker will be splitting.
    """
    return {
        "rows": len(rows),
        "columns": max((len(row.split(_CELL_SEPARATOR)) for row in rows), default=0),
        "header_rows": 1 if table.first_row and rows else 0,
    }


def _shape_items(
    shapes: Iterable[BaseShape],
    size: tuple[int, int] | None,
    title_id: int | None,
    *,
    grouped: bool,
) -> Iterator[_SlideItem]:
    """Yield an item per text-bearing shape, in the order the slide stores them."""
    for shape in shapes:
        if isinstance(shape, GroupShape):
            # A grouped shape's offsets are in the group's child coordinate space, and mapping
            # them onto the slide needs the group's chOff/chExt transform, which python-pptx
            # does not resolve. Its text is still content, so it is emitted page-level rather
            # than dropped or given a box measured in the wrong space.
            yield from _shape_items(shape.shapes, size, title_id, grouped=True)
            continue

        rect = None if grouped else _shape_rect(shape, size)

        if isinstance(shape, GraphicFrame) and shape.has_table:
            rows = _render_table(shape.table)
            text = _ROW_SEPARATOR.join(rows)
            if text.strip():
                yield _SlideItem(
                    kind=BlockKind.TABLE,
                    text=text,
                    rect=rect,
                    metadata=_table_metadata(shape.table, rows),
                )
            continue

        if not isinstance(shape, Shape):
            # A picture, a chart or a media object. With optical character recognition out of
            # scope there is no text in it, and python-pptx exposes no alternative text, so
            # there is nothing to index — which docs/parsing.md §4.1 states is the honest
            # outcome rather than a gap to fill with the shape's name.
            continue

        frame = shape.text_frame
        text = frame.text.strip()
        if not text:
            continue
        if shape.shape_id == title_id:
            yield _SlideItem(kind=BlockKind.HEADING, text=text, rect=rect)
            continue
        levels = [paragraph.level for paragraph in frame.paragraphs]
        # An indent level is the only list signal python-pptx resolves: bullet characters live
        # in inherited list styles it does not read. So a frame with an indented paragraph is a
        # list, and a flat one is prose rather than a guessed-at bulleted list.
        metadata: Metadata = {}
        if any(level > 0 for level in levels):
            metadata = {"indent_levels": list(levels)}
            yield _SlideItem(kind=BlockKind.LIST, text=text, rect=rect, metadata=metadata)
            continue
        yield _SlideItem(kind=BlockKind.PROSE, text=text, rect=rect, metadata=metadata)


def _notes_text(slide: Slide) -> str:
    """The slide's speaker notes, or ``""``.

    ``has_notes_slide`` is checked first because reading ``notes_slide`` creates one, which
    would make an inspection change the thing being inspected.
    """
    if not slide.has_notes_slide:
        return ""
    frame = slide.notes_slide.notes_text_frame
    return "" if frame is None else frame.text.strip()


def _slide_items(
    slide: Slide, size: tuple[int, int] | None, config: SlidesConfig
) -> list[_SlideItem]:
    """Every block-to-be on one slide, shapes in z-order and then the speaker notes."""
    title = slide.shapes.title
    drafts = list(
        _shape_items(slide.shapes, size, None if title is None else title.shape_id, grouped=False)
    )

    occurrences = Counter(item.rect for item in drafts if item.rect is not None)
    items = [
        item if item.rect is None or occurrences[item.rect] == 1 else replace(item, rect=None)
        for item in drafts
    ]

    if config.include_speaker_notes:
        notes = _notes_text(slide)
        if notes:
            items.append(
                _SlideItem(
                    kind=BlockKind.PROSE, text=notes, rect=None, metadata={"speaker_notes": True}
                )
            )
    return items


def _heading_path(slide: Slide) -> tuple[str, ...]:
    """The slide title as a heading path, or empty when the slide has no title.

    Empty rather than "Slide 4": the path feeds the breadcrumb, which feeds the embedding, and
    a positional label there is a signal about nothing.
    """
    title = slide.shapes.title
    if title is None:
        return ()
    text = title.text_frame.text.strip()
    return (text,) if text else ()


def _same_rect(left: Rect, right: Rect) -> bool:
    return all(
        abs(a - b) <= _COORDINATE_TOLERANCE
        for a, b in (
            (left.x0, right.x0),
            (left.y0, right.y0),
            (left.x1, right.x1),
            (left.y1, right.y1),
        )
    )


class SlidesParser:
    """Parses a `.pptx` into per-slide blocks anchored by slide number and shape box.

    ``max_pagelevel_ratio`` is 0.20 (``docs/parsing.md`` §3.4). It is not slack for
    convenience: speaker notes, grouped shapes and shapes with no reported geometry are all
    legitimately page-level, and the ceiling is what makes it visible if box extraction stops
    working — the citation still names the right slide, so nothing else would say.
    """

    media_types = frozenset({MEDIA_TYPE})
    profile = ParserProfile(name="slides", max_unlocated_ratio=0.0, max_pagelevel_ratio=0.20)

    def __init__(self, config: SlidesConfig) -> None:
        self._config = config

    async def parse(self, raw: RawDocument) -> AsyncIterator[ParsedBlock]:
        """Yield blocks slide by slide, shapes in z-order, notes last on their slide."""
        presentation = _open(raw)
        size = _slide_size(presentation)
        for page, slide in enumerate(presentation.slides, start=1):
            heading_path = _heading_path(slide)
            for item in _slide_items(slide, size, self._config):
                yield ParsedBlock(
                    kind=item.kind,
                    text=item.text,
                    anchor=PageAnchor(page=page, rects=() if item.rect is None else (item.rect,)),
                    heading_path=heading_path,
                    metadata={"slide_id": slide.slide_id, **item.metadata},
                )

    async def resolve(self, anchor: Anchor, raw: RawDocument) -> str | None:
        """Return the text ``anchor`` addresses, re-derived from ``raw``.

        With rectangles, that is the text of the shapes those boxes cover — several, when the
        chunker merged blocks from two shapes into one chunk. Without them, the whole slide,
        which is what a page-level anchor honestly claims.
        """
        if not isinstance(anchor, PageAnchor):
            return None
        presentation = _open(raw)
        slides = list(presentation.slides)
        if anchor.page > len(slides):
            return None
        items = _slide_items(slides[anchor.page - 1], _slide_size(presentation), self._config)
        if not anchor.rects:
            return _ROW_SEPARATOR.join(item.text for item in items) or None
        matched = [
            item.text
            for item in items
            if item.rect is not None
            and any(_same_rect(item.rect, wanted) for wanted in anchor.rects)
        ]
        return _ROW_SEPARATOR.join(matched) or None


def _open(raw: RawDocument) -> Presentation:
    """Open the package, declining anything that is not a readable PresentationML one."""
    try:
        return pptx.Presentation(io.BytesIO(raw.as_bytes()))
    # lxml's XMLSyntaxError subclasses SyntaxError and is caught by that name rather than by
    # importing lxml, so declining a truncated ppt/presentation.xml does not tie this module
    # to python-pptx's choice of XML backend.
    except (PackageNotFoundError, zipfile.BadZipFile, KeyError, SyntaxError, ValueError) as exc:
        msg = (
            f"{raw.uri}: not a readable .pptx package ({type(exc).__name__}: {exc}). Expected "
            f"an OOXML PresentationML package containing ppt/presentation.xml. Re-export it "
            f"from a presentation editor, or route this media type to a different parser"
        )
        raise ParseError(msg) from exc
